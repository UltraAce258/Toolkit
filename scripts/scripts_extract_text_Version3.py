#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
[display-name-zh] 通用文本提取器
[display-name-en] Universal Text Extractor

功能:
  批量提取文档、图片、音频和视频中的文本内容。音视频将通过 ffmpeg 转为音频后调用 Whisper 生成文本与字幕。

---
兼容性:
- 文件格式:
  - 文档: .pdf, .docx, .pptx, .doc (旧版Word, 仅限Windows), .json
  - 图片: .jpg, .jpeg, .png, .bmp, .webp, .heic, .heif, .raw, .cr2, .nef, .arw, .dng
  - 音频/视频: .mp3, .wav, .m4a, .aac, .flac, .ogg, .mp4, .mkv, .mov, .avi, .webm
- 平台: 跨平台（.doc 为 Windows 特性）
---
更新日志:
  - v5.0 (2026-04-15): 修复 OCR CUDA 路径；集成音视频字幕提取（ffmpeg + faster-whisper）；重构提取分发逻辑。
~~~
Function:
  Batch extracts text from documents, images, audio and video. For media files, it converts input to audio using ffmpeg and generates both transcript text and subtitles via Whisper.

---
Compatibility:
- File Formats:
  - Documents: .pdf, .docx, .pptx, .doc (legacy Word, Windows only), .json
  - Images: .jpg, .jpeg, .png, .bmp, .webp, .heic, .heif, .raw, .cr2, .nef, .arw, .dng
  - Audio/Video: .mp3, .wav, .m4a, .aac, .flac, .ogg, .mp4, .mkv, .mov, .avi, .webm
- Platform: Cross-platform (.doc is Windows-only)
---
Changelog:
  - v5.0 (2026-04-15): Fixed OCR CUDA path; integrated media subtitle extraction (ffmpeg + faster-whisper); refactored extractor dispatch logic.
"""

import os
import sys
import argparse
import platform
import subprocess
import importlib.util
import shutil
import tempfile
from pathlib import Path
from dataclasses import dataclass

# --- Internationalization (i18n) Setup ---
MESSAGES = {
    'zh': {
        "init": "--- 通用文本提取器 v5.0 启动 ---",
        "dep_checking": "--- 正在检查依赖库 ---",
        "dep_missing": "提示: 未找到 '{module}'，正在尝试自动安装 '{pkg}'...",
        "dep_success": "成功: '{pkg}' 已安装。",
        "dep_fail": "错误: 自动安装 '{pkg}' 失败。请手动运行 'pip install {pkg}'。",
        "dep_warn_ocr": "警告: OCR引擎 'easyocr' 安装失败。图片提取功能将不可用。",
        "dep_warn_media": "警告: Whisper 依赖安装失败。音视频提取将不可用。",
        "search_mode": "搜索模式: 自动 (文件直接处理，文件夹递归搜索)",
        "mode_sort_by": "排序方式: {sort_by}",
        "mode_ocr_device": "OCR 设备模式: {device}",
        "mode_subtitle_mode": "字幕提取模式: {mode}",
        "files_found": "\n发现 {count} 个待处理文件:",
        "files_none": "\n未找到任何支持的文档、图片或音视频文件。",
        "ocr_init": "首次运行，正在初始化OCR引擎 (可能需要下载模型)...",
        "ocr_ready": "OCR引擎准备就绪。",
        "ocr_using": "OCR 当前运行设备: {device}",
        "ocr_cuda_unavailable": "检测到 CUDA 不可用，OCR 将自动回退 CPU。建议安装 GPU 版 PyTorch 并更新驱动。",
        "ffmpeg_missing": "错误: 未找到 ffmpeg。音视频字幕提取将被跳过。",
        "ffmpeg_guide": "安装建议: Windows(winget install ffmpeg) / macOS(brew install ffmpeg) / Ubuntu(sudo apt install ffmpeg)",
        "media_extracting": "  [媒体] 正在通过 ffmpeg 提取音频流...",
        "media_transcribing": "  [媒体] 正在调用 Whisper 转写...",
        "media_srt_saved": "  [成功] -> 字幕已保存到: {path}",
        "media_fail": "  [失败] -> 音视频提取失败: {e}",
        "media_skip_no_engine": "  [跳过] -> 当前环境未就绪，无法处理音视频文件。",
        "processing": "\n[处理中 {i}/{total}] -> {filename}",
        "success_save": "  [成功] -> 已保存到: {path}",
        "failure_write": "  [失败] -> 写入文件时出错: {e}",
        "failure_process": "  [失败] -> 处理文件时出错: {e}",
        "failure_read_img": "  [失败] -> 读取图片失败: {e}",
        "unsupported": "  跳过不支持的文件。",
        "doc_only_windows": "错误: .doc文件仅在Windows系统且安装了Word后才能处理。",
        "output_dir_creating": "  创建输出目录: {path}",
        "output_dir_fail": "  错误: 创建输出目录失败: {e}",
        "all_done": "\n--- 所有任务已完成 ---",
    },
    'en': {
        "init": "--- Universal Text Extractor v5.0 Started ---",
        "dep_checking": "--- Checking Dependencies ---",
        "dep_missing": "Info: '{module}' not found. Attempting to auto-install '{pkg}'...",
        "dep_success": "Success: '{pkg}' has been installed.",
        "dep_fail": "Error: Auto-install of '{pkg}' failed. Please run 'pip install {pkg}' manually.",
        "dep_warn_ocr": "Warning: OCR engine 'easyocr' failed to install. Image extraction will be unavailable.",
        "dep_warn_media": "Warning: Whisper dependencies failed to install. Audio/video extraction will be unavailable.",
        "search_mode": "Search Mode: Automatic (Files processed directly, folders searched recursively)",
        "mode_sort_by": "Sorting by: {sort_by}",
        "mode_ocr_device": "OCR device mode: {device}",
        "mode_subtitle_mode": "Subtitle mode: {mode}",
        "files_found": "\nFound {count} files to process:",
        "files_none": "\nNo supported document, image, audio, or video files found.",
        "ocr_init": "First run, initializing OCR engine (may download models)...",
        "ocr_ready": "OCR engine is ready.",
        "ocr_using": "OCR runtime device: {device}",
        "ocr_cuda_unavailable": "CUDA is unavailable. OCR falls back to CPU. Recommendation: install a GPU-enabled PyTorch build and update your driver.",
        "ffmpeg_missing": "Error: ffmpeg not found. Audio/video subtitle extraction will be skipped.",
        "ffmpeg_guide": "Install guide: Windows(winget install ffmpeg) / macOS(brew install ffmpeg) / Ubuntu(sudo apt install ffmpeg)",
        "media_extracting": "  [MEDIA] Extracting audio stream via ffmpeg...",
        "media_transcribing": "  [MEDIA] Running Whisper transcription...",
        "media_srt_saved": "  [SUCCESS] -> Subtitle saved to: {path}",
        "media_fail": "  [FAILURE] -> Media extraction failed: {e}",
        "media_skip_no_engine": "  [SKIP] -> Runtime environment is not ready for media files.",
        "processing": "\n[Processing {i}/{total}] -> {filename}",
        "success_save": "  [SUCCESS] -> Saved to: {path}",
        "failure_write": "  [FAILURE] -> Error writing file: {e}",
        "failure_process": "  [FAILURE] -> Error processing file: {e}",
        "failure_read_img": "  [FAILURE] -> Failed to read image: {e}",
        "unsupported": "  Skipping unsupported file.",
        "doc_only_windows": "Error: .doc files can only be processed on Windows with MS Word installed.",
        "output_dir_creating": "  Creating output directory: {path}",
        "output_dir_fail": "  Error: Failed to create output directory: {e}",
        "all_done": "\n--- All tasks completed ---",
    }
}

def T(key, lang='en', **kwargs):
    return MESSAGES.get(lang, MESSAGES['en']).get(key, key).format(**kwargs)


DOC_EXTS = ['.pdf', '.docx', '.pptx', '.doc', '.json']
IMG_EXTS = ['.jpg', '.jpeg', '.png', '.bmp', '.webp', '.heic', '.heif', '.raw', '.cr2', '.nef', '.arw', '.dng']
MEDIA_EXTS = ['.mp3', '.wav', '.m4a', '.aac', '.flac', '.ogg', '.mp4', '.mkv', '.mov', '.avi', '.webm']
SUPPORTED_EXTS = DOC_EXTS + IMG_EXTS + MEDIA_EXTS

# --- Dependency Management ---
def check_and_install(pkg_name, module_name=None, lang='en', extra_args=None):
    if module_name is None: module_name = pkg_name
    if importlib.util.find_spec(module_name): return True
    print(T("dep_missing", lang, pkg=pkg_name, module=module_name))
    try:
        command = [sys.executable, "-m", "pip", "install", pkg_name]
        if extra_args: command.extend(extra_args)
        subprocess.check_call(command, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        print(T("dep_success", lang, pkg=pkg_name))
        importlib.invalidate_caches()
        return True
    except subprocess.CalledProcessError:
        print(T("dep_fail", lang, pkg=pkg_name))
        return False

def setup_dependencies(lang, need_docs=True, need_images=False, need_media=False):
    print(T("dep_checking", lang))
    if need_docs:
        check_and_install('python-docx', 'docx', lang)
        check_and_install('python-pptx', 'pptx', lang)
        check_and_install('PyMuPDF', 'fitz', lang)
    if platform.system() == "Windows":
        check_and_install('pywin32', 'win32com', lang)
    check_and_install('natsort', 'natsort', lang)
    if need_images:
        check_and_install('pillow-heif', 'pillow_heif', lang)
        check_and_install('rawpy', 'rawpy', lang)
        if not check_and_install('easyocr', lang=lang):
            print(T("dep_warn_ocr", lang))
    if need_media:
        ok_fw = check_and_install('faster-whisper', 'faster_whisper', lang)
        if not ok_fw:
            print(T("dep_warn_media", lang))

# --- Extraction Functions ---
def extract_text_from_docx(file_path):
    import docx
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    import pptx
    prs = pptx.Presentation(file_path)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def extract_text_from_pdf(file_path):
    import fitz
    with fitz.open(file_path) as doc:
        return "".join([page.get_text() for page in doc])

def extract_text_from_doc(file_path, lang):
    if platform.system() != "Windows": return T("doc_only_windows", lang)
    try:
        import win32com.client, pythoncom
        pythoncom.CoInitialize()
        word_app = win32com.client.Dispatch("Word.Application")
        word_app.Visible = False
        doc = word_app.Documents.Open(os.path.abspath(file_path))
        text = doc.Content.Text
        doc.Close()
        word_app.Quit()
        pythoncom.CoUninitialize()
        return text
    except Exception: return T("doc_only_windows", lang)
    
def extract_text_from_json(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_image(file_path, ocr_reader, lang):
    import numpy as np
    from PIL import Image
    try:
        ext = Path(file_path).suffix.lower()
        if ext in ['.heic', '.heif']:
            import pillow_heif
            heif_file = pillow_heif.read_heif(file_path)
            image = Image.frombytes(heif_file.mode, heif_file.size, heif_file.data, "raw")
        elif ext in ['.raw', '.cr2', '.nef', '.arw', '.dng']:
            import rawpy
            with rawpy.imread(file_path) as raw: image = Image.fromarray(raw.postprocess())
        else: image = Image.open(file_path)
        image_np = np.array(image.convert('RGB'))
        result = ocr_reader.readtext(image_np, detail=0, paragraph=True)
        return "\n".join(result)
    except Exception as e: raise IOError(T("failure_read_img", lang, e=e)) from e


@dataclass
class SegmentLite:
    start: float
    end: float
    text: str
    no_speech_prob: float
    lang: str


def _torch_cuda_available():
    try:
        import torch
        return bool(torch.cuda.is_available())
    except Exception:
        return False


def resolve_runtime_device(preferred, lang):
    if preferred == "cpu":
        return "cpu"
    if preferred == "cuda":
        if _torch_cuda_available():
            return "cuda"
        print(T("ocr_cuda_unavailable", lang))
        return "cpu"
    return "cuda" if _torch_cuda_available() else "cpu"


def get_ocr_reader(lang, ocr_device):
    try:
        import easyocr
        resolved = resolve_runtime_device(ocr_device, lang)
        print(T("ocr_init", lang))
        reader = easyocr.Reader(['ch_sim', 'en'], gpu=(resolved == "cuda"))
        print(T("ocr_ready", lang))
        print(T("ocr_using", lang, device=resolved))
        return reader
    except Exception:
        return None


def ensure_ffmpeg(lang):
    if shutil.which("ffmpeg"):
        return True
    print(T("ffmpeg_missing", lang))
    print(T("ffmpeg_guide", lang))
    return False


def format_srt_time(seconds):
    m, s = divmod(seconds, 60)
    h, m = divmod(m, 60)
    return f"{int(h):02d}:{int(m):02d}:{int(s):02d},{int((seconds - int(s)) * 1000):03d}"


def segments_to_srt(segments):
    rows = []
    for i, seg in enumerate(segments, 1):
        rows.append(f"{i}\n{format_srt_time(seg.start)} --> {format_srt_time(seg.end)}\n{seg.text.strip()}\n")
    return "\n".join(rows).strip() + "\n"


def _extract_audio_for_whisper(media_file, lang):
    print(T("media_extracting", lang))
    temp_dir = tempfile.mkdtemp(prefix="toolkit_media_")
    wav_path = Path(temp_dir) / f"{Path(media_file).stem}_extracted.wav"
    command = [
        "ffmpeg", "-y", "-i", str(media_file), "-vn", "-ac", "1", "-ar", "16000", str(wav_path)
    ]
    subprocess.run(command, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
    return wav_path


def _run_transcribe(model, audio_path, mode):
    kwargs = {"beam_size": 5, "no_speech_threshold": 0.6}
    if mode in ("en", "zh"):
        kwargs["language"] = mode
    segments_generator, _ = model.transcribe(str(audio_path), **kwargs)
    return [
        SegmentLite(
            start=float(seg.start),
            end=float(seg.end),
            text=str(seg.text).strip(),
            no_speech_prob=float(getattr(seg, "no_speech_prob", 1.0)),
            lang=mode if mode in ("en", "zh") else "auto",
        )
        for seg in segments_generator
    ]


def intelligent_merge(en_segments, zh_segments):
    final_segments = []
    zh_pool = list(zh_segments)
    for en_seg in en_segments:
        overlapping = [z for z in zh_pool if max(en_seg.start, z.start) < min(en_seg.end, z.end)]
        if not overlapping:
            final_segments.append(en_seg)
            continue
        avg_zh_prob = sum(seg.no_speech_prob for seg in overlapping) / len(overlapping)
        if en_seg.no_speech_prob < avg_zh_prob:
            final_segments.append(en_seg)
        else:
            final_segments.extend(overlapping)
        for z in overlapping:
            if z in zh_pool:
                zh_pool.remove(z)
    final_segments.extend(zh_pool)
    final_segments.sort(key=lambda s: s.start)
    return final_segments


def extract_from_media(file_path, whisper_model, whisper_device, subtitle_mode, lang):
    from faster_whisper import WhisperModel
    resolved_device = resolve_runtime_device(whisper_device, lang)
    compute_type = "int8_float16" if resolved_device == "cuda" else "int8"
    model = WhisperModel(whisper_model, device=resolved_device, compute_type=compute_type)
    audio_path = _extract_audio_for_whisper(file_path, lang)
    print(T("media_transcribing", lang))
    try:
        if subtitle_mode == "merge":
            en_segments = _run_transcribe(model, audio_path, "en")
            zh_segments = _run_transcribe(model, audio_path, "zh")
            merged = intelligent_merge(en_segments, zh_segments)
            return "\n".join(seg.text for seg in merged), segments_to_srt(merged)
        segments = _run_transcribe(model, audio_path, subtitle_mode)
        return "\n".join(seg.text for seg in segments), segments_to_srt(segments)
    finally:
        try:
            if audio_path.exists():
                audio_path.unlink()
            parent = audio_path.parent
            if parent.exists():
                parent.rmdir()
        except Exception:
            pass

# --- Main Logic ---
def get_all_supported_files(paths, sort_by):
    files_to_process = set()
    for path in paths:
        p = Path(path)
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS:
            files_to_process.add(str(p))
        elif p.is_dir():
            for ext in SUPPORTED_EXTS:
                files_to_process.update(str(f) for f in p.rglob(f"*{ext}"))
    file_list = list(files_to_process)
    if sort_by == 'date': file_list.sort(key=os.path.getmtime)
    else:
        try:
            from natsort import natsorted
            file_list = natsorted(file_list)
        except Exception:
            file_list.sort()
    return file_list

def main():
    parser = argparse.ArgumentParser(description="Universal Text Extractor for documents and images.")
    parser.add_argument('files', nargs='*', help="Paths to files or folders from the GUI.")
    parser.add_argument('--gui-mode', action='store_true', help=argparse.SUPPRESS)
    parser.add_argument('--lang', type=str, default='en', choices=['zh', 'en'], help=argparse.SUPPRESS)
    parser.add_argument('--sort-by', type=str, default="name", choices=["name", "date"], help="Sort files by 'name' (natural) or 'date' (modification time).")
    parser.add_argument('--ocr-device', type=str, default='auto', choices=['auto', 'cpu', 'cuda'], help="OCR device: auto/cpu/cuda.")
    parser.add_argument('--whisper-model', type=str, default='small', help="Whisper model name for media subtitle extraction.")
    parser.add_argument('--whisper-device', type=str, default='auto', choices=['auto', 'cpu', 'cuda'], help="Whisper runtime device: auto/cpu/cuda.")
    parser.add_argument('--subtitle-mode', type=str, default='auto', choices=['auto', 'merge', 'en', 'zh'], help="Subtitle extraction mode for media files.")
    args = parser.parse_args()
    lang = args.lang

    print(T("init", lang))
    print(T("search_mode", lang))
    print(T("mode_sort_by", lang, sort_by=args.sort_by))
    print(T("mode_ocr_device", lang, device=args.ocr_device))
    print(T("mode_subtitle_mode", lang, mode=args.subtitle_mode))

    files_to_process = get_all_supported_files(args.files, args.sort_by)

    if not files_to_process:
        print(T("files_none", lang))
        return

    need_docs = any(Path(f).suffix.lower() in DOC_EXTS for f in files_to_process)
    need_images = any(Path(f).suffix.lower() in IMG_EXTS for f in files_to_process)
    need_media = any(Path(f).suffix.lower() in MEDIA_EXTS for f in files_to_process)
    setup_dependencies(lang, need_docs=need_docs, need_images=need_images, need_media=need_media)

    print(T("files_found", lang, count=len(files_to_process)))
    for f in files_to_process: print(f"- {f}")

    # ======================================================================
    # vvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvv
    # (MODIFIED) Context-Aware Output Directory Logic
    # ======================================================================
    output_base_dir = None
    if len(files_to_process) > 1:
        parent_dirs = {Path(f).parent for f in files_to_process}
        if len(parent_dirs) == 1:
            # Batch mode, same source directory
            output_base_dir = list(parent_dirs)[0] / "提取的文本_Extracted_Text"
        else:
            # Batch mode, multiple source directories
            output_base_dir = Path.cwd() / "提取的文本_Extracted_Text"
    # If single file mode, output_base_dir remains None
    # ======================================================================
    # ^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^

    ocr_reader = get_ocr_reader(lang, args.ocr_device) if need_images else None
    media_runtime_ok = ensure_ffmpeg(lang) if need_media else True

    for i, file_path in enumerate(files_to_process):
        p = Path(file_path)
        print(T("processing", lang, i=i+1, total=len(files_to_process), filename=p.name))
        
        # Determine final output directory for this specific file
        if output_base_dir:
            output_dir = output_base_dir
        else:
            # Single file mode
            output_dir = p.parent

        # Create output directory if it doesn't exist
        if not output_dir.exists():
            try:
                print(T("output_dir_creating", lang, path=output_dir))
                output_dir.mkdir(parents=True)
            except OSError as e:
                print(T("output_dir_fail", lang, e=e))
                continue # Skip this file if its output dir can't be created

        output_path = output_dir / f"{p.stem}_extracted.txt"
        subtitle_path = output_dir / f"{p.stem}_extracted.srt"
        text_content = ""
        
        try:
            ext = p.suffix.lower()
            if ext in ['.pdf']: text_content = extract_text_from_pdf(p)
            elif ext in ['.docx']: text_content = extract_text_from_docx(p)
            elif ext in ['.pptx']: text_content = extract_text_from_pptx(p)
            elif ext in ['.doc']: text_content = extract_text_from_doc(p, lang)
            elif ext in ['.json']: text_content = extract_text_from_json(p)
            elif ext in IMG_EXTS:
                if ocr_reader: text_content = extract_text_from_image(p, ocr_reader, lang)
                else: continue
            elif ext in MEDIA_EXTS:
                if not media_runtime_ok:
                    print(T("media_skip_no_engine", lang))
                    continue
                try:
                    text_content, srt_content = extract_from_media(
                        p, args.whisper_model, args.whisper_device, args.subtitle_mode, lang
                    )
                    with open(subtitle_path, 'w', encoding='utf-8') as sf:
                        sf.write(srt_content)
                    print(T("media_srt_saved", lang, path=subtitle_path))
                except Exception as e:
                    print(T("media_fail", lang, e=e))
                    continue
            else:
                print(T("unsupported", lang)); continue

            with open(output_path, 'w', encoding='utf-8') as f: f.write(text_content)
            print(T("success_save", lang, path=output_path))
        except Exception as e:
            print(T("failure_process", lang, e=e))

    print(T("all_done", lang))

if __name__ == "__main__":
    main()
